#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_builder.core
-----------------
Core module for PPTX Builder: Create PowerPoint presentations from PDFs and images.

__version__ = "0.1.0"

Features:
    - Prompts for source folder and output filename.
    - Lets you choose common slide sizes (16:9, 4:3, Letter, A4, Legal, Tabloid).
    - Two placement modes:
        (1) Fit whole image (no crop, may show background bars),
        (2) Crop-to-fill (cover slide fully, cropping as needed; no distortion).
    - One image per slide, centered, sorted by filename (case-insensitive).

Notes:
    - Uses python-pptx (and Pillow) under the hood.
    - Slide background shows through when using "Fit whole image".
    - No part of the image is stretched; scaling is always proportional.
"""

import sys
import logging
from pathlib import Path
from typing import Tuple, List

from pptx import Presentation
from pptx.util import Inches, Emu
from tqdm import tqdm

# Set up module logger
logger = logging.getLogger(__name__)

# -----------------------------
# Slide size presets (in inches)
# -----------------------------
SLIDE_SIZES = {
    "1": ('16:9 (13.33" x 7.5")', 13.3333333333, 7.5),
    "2": ('4:3  (10" x 7.5")', 10.0, 7.5),
    "3": ('Letter  (11" x 8.5")', 11.0, 8.5),
    "4": ('A4      (11.69" x 8.27")', 11.69, 8.27),
    "5": ('Legal   (14" x 8.5")', 14.0, 8.5),
    "6": ('Tabloid (17" x 11")', 17.0, 11.0),
}

# -----------------------------
# File extensions we will accept
# -----------------------------
ALLOWED_EXTS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".tif",
    ".tiff",
    ".webp",
    ".bmp",
    ".gif",
    ".ico",
    ".heic",
    ".heif",
}


def prompt_input_path() -> Path:
    """Ask the user for a path to a PDF file or a folder of images."""
    while True:
        path_str = input("Enter the path to a PDF file or a folder of images: ").strip()
        path = Path(path_str).expanduser().resolve()
        if path.exists():
            return path
        print("✗ That path does not exist. Please try again.\n")


def prompt_output_name(default_name: str = "output") -> str:
    """Ask the user for the output file name (no extension needed)."""
    out = input(f"Enter output filename (without extension) [{default_name}]: ").strip()
    if not out:
        out = default_name
    # Ensure .pptx extension
    if not out.lower().endswith(".pptx"):
        out = out + ".pptx"
    return out


def prompt_slide_size() -> Tuple[float, float]:
    """Display size options and return (width_in_inches, height_in_inches)."""
    print("\nChoose slide size:")
    for key, (label, w, h) in SLIDE_SIZES.items():
        print(f"  {key}) {label}")
    while True:
        choice = input("Enter number (1-6): ").strip()
        if choice in SLIDE_SIZES:
            _, w, h = SLIDE_SIZES[choice]
            return w, h
        print("✗ Invalid choice. Please enter a number from the list.\n")


def prompt_fit_mode() -> str:
    """
    Ask how images should be placed:
        1) Fit whole image (no cropping)  -> 'fit'
        2) Crop to fill (cover, may crop) -> 'fill'
    """
    print("\nHow should images be placed?")
    print("  1) Fit whole image (no cropping; background may show)")
    print("  2) Crop to fill (no whitespace; edges may be trimmed)")
    while True:
        choice = input("Enter 1 or 2: ").strip()
        if choice == "1":
            return "fit"
        if choice == "2":
            return "fill"
        print("✗ Invalid choice. Please enter 1 or 2.\n")


def list_images(folder: Path) -> List[Path]:
    """Return sorted list of image files with allowed extensions (case-insensitive)."""
    files = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in ALLOWED_EXTS]
    # Sort case-insensitively by filename
    files.sort(key=lambda p: p.name.lower())
    return files


def emu_to_float_inches(emu: Emu) -> float:
    """Convert EMU to inches (pptx.util.Inches wraps conversion, but we need a float)."""
    # 1 inch = 914400 EMU
    return float(emu) / 914400.0


def place_picture_fit(slide, img_path: Path, slide_w_emu: int, slide_h_emu: int):
    """
    Place the image on the slide using 'contain' behavior:
        - Scale proportionally so the entire image is visible (no cropping).
        - Center the image; background may show (letterbox/pillarbox).
    Implementation detail:
        - Insert at natural size first to read pic.width/height (requires Pillow via python-pptx).
        - Compute scale ratio and then set final size and position.
    """
    pic = slide.shapes.add_picture(str(img_path), left=0, top=0)  # natural size first
    img_w = float(pic.width)
    img_h = float(pic.height)
    sw = float(slide_w_emu)
    sh = float(slide_h_emu)

    # Scale to "contain": use the smaller ratio
    scale = min(sw / img_w, sh / img_h)
    new_w = img_w * scale
    new_h = img_h * scale

    # Center
    left = (sw - new_w) / 2.0
    top = (sh - new_h) / 2.0

    pic.left = int(left)
    pic.top = int(top)
    pic.width = int(new_w)
    pic.height = int(new_h)


def place_picture_fill(slide, img_path: Path, slide_w_emu: int, slide_h_emu: int):
    """
    Place the image on the slide using 'cover' behavior:
        - Scale proportionally so the slide is fully covered.
        - May crop the image (overflow outside slide bounds is not visible).
        - Center the image.
    """
    pic = slide.shapes.add_picture(str(img_path), left=0, top=0)  # natural size first
    img_w = float(pic.width)
    img_h = float(pic.height)
    sw = float(slide_w_emu)
    sh = float(slide_h_emu)

    # Scale to "cover": use the larger ratio
    scale = max(sw / img_w, sh / img_h)
    new_w = img_w * scale
    new_h = img_h * scale

    # Center (image may overflow the slide; that's fine)
    left = (sw - new_w) / 2.0
    top = (sh - new_h) / 2.0

    pic.left = int(left)
    pic.top = int(top)
    pic.width = int(new_w)
    pic.height = int(new_h)


def confirm_overwrite(path: Path, quiet: bool = False, force: bool = False) -> bool:
    if not path.exists():
        return True
    if quiet or force:
        return True
    reply = input(f"⚠️  File exists: {path.name}. Overwrite? [y/N]: ").strip().lower()
    return reply == "y"


def build_presentation(
    images: List[Path],
    output_path: Path,
    slide_width_in: float,
    slide_height_in: float,
    mode: str,
    show_progress: bool = False,
) -> None:
    """Create the PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    # Cache slide size in EMU for math (we just set these, so they're not None)
    sw_emu = int(prs.slide_width)  # type: ignore[arg-type]
    sh_emu = int(prs.slide_height)  # type: ignore[arg-type]

    # Create slides with optional progress bar
    image_iter = tqdm(images, desc="Building slides", unit="slide") if show_progress else images
    for img in image_iter:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        if mode == "fit":
            place_picture_fit(slide, img, sw_emu, sh_emu)
        else:
            place_picture_fill(slide, img, sw_emu, sh_emu)

    if show_progress:
        print("Saving presentation...")
    prs.save(str(output_path))


# ===[ SECTION: CLI ARGUMENTS ]=====================================

import argparse  # noqa: E402


def parse_cli_args():
    """Parse command-line arguments for batch, recursive, or quiet runs."""
    parser = argparse.ArgumentParser(
        prog="pptx-builder",
        description="Build PowerPoint (.pptx) files from PDFs or image folders.",
    )

    parser.add_argument(
        "-i",
        "--input",
        nargs="+",
        metavar="PATH",
        help="Path(s) to one or more PDFs, images, or folders to process.",
    )

    parser.add_argument(
        "-r",
        "--recursive",
        action="store_true",
        help="Recurse into subfolders when processing folders.",
    )

    parser.add_argument(
        "--dpi", type=int, default=300, help="DPI for PDF rendering (default: 300)."
    )

    parser.add_argument(
        "--quiet",
        action="store_true",
        help="Suppress interactive prompts and non-critical output.",
    )

    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable verbose debug logging output.",
    )

    parser.add_argument(
        "--force",
        action="store_true",
        help="Overwrite existing PPTX files without confirmation",
    )

    parser.add_argument(
        "-o",
        "--output",
        metavar="NAME",
        help="Output filename (with or without .pptx extension). Only valid for single input.",
    )

    return parser.parse_args()


# ===[ SECTION: INPUT HANDLING ]====================================

from pdf2image import convert_from_path  # noqa: E402


def detect_input_type(path: Path) -> str:
    """Return 'pdf', 'folder', or 'unknown' based on the given path."""
    if path.is_file() and path.suffix.lower() == ".pdf":
        return "pdf"
    if path.is_dir():
        # check if folder contains images
        imgs = list_images(path)
        if imgs:
            return "folder"
    return "unknown"


def convert_pdf_to_images(pdf_path: Path, dpi: int) -> List[Path]:
    """Convert PDF pages to temporary PNG files."""
    import tempfile  # noqa: E402

    logger.debug(f"Starting PDF conversion: {pdf_path}")
    logger.debug(f"PDF exists: {pdf_path.exists()}")
    logger.debug(f"PDF size: {pdf_path.stat().st_size if pdf_path.exists() else 'N/A'}")

    temp_dir = Path(tempfile.mkdtemp(prefix="pptx_pdf_"))
    logger.debug(f"Created temp dir: {temp_dir}")

    try:
        logger.debug(f"Calling convert_from_path with dpi={dpi}")
        pages = convert_from_path(pdf_path.as_posix(), dpi=dpi)
        logger.debug(f"Got {len(pages)} pages from PDF")
        out_paths = []

        # Convert pages with progress bar
        page_iter = tqdm(
            enumerate(pages, start=1),
            total=len(pages),
            desc="Converting PDF pages",
            unit="page",
        )
        for i, page in page_iter:
            out_path = temp_dir / f"page_{i:04d}.png"
            page.save(out_path, "PNG")
            out_paths.append(out_path)
            logger.debug(f"Saved page {i} to {out_path}")

        logger.debug(f"Successfully converted {len(out_paths)} pages")
        return out_paths
    except Exception as e:
        # Clean up temp directory on failure
        import shutil
        import traceback

        logger.error(f"PDF conversion failed: {e}")
        logger.debug(f"Traceback:\n{traceback.format_exc()}")
        if temp_dir.exists():
            shutil.rmtree(temp_dir)
        raise RuntimeError(f"Failed to convert PDF: {e}")


def pdf_first_page_size_inches(pdf_path: Path) -> Tuple[float, float]:
    """
    Return (width_in, height_in) for the first page of a PDF.
    Uses 72 PDF points per inch via PyMuPDF (fitz).
    """
    import fitz  # PyMuPDF

    with fitz.open(pdf_path) as doc:
        if doc.page_count == 0:
            # safe fallback to 16:9 if something is odd
            return 13.3333, 7.5
        p0 = doc[0]
        w_in = p0.rect.width / 72.0
        h_in = p0.rect.height / 72.0
        return (w_in, h_in)


def process_folder(folder: Path, recursive: bool, dpi: int, quiet: bool) -> None:
    """Process all PDFs and/or images in a folder into PPTX files."""
    pdfs = sorted(folder.glob("*.pdf"))
    imgs = [p for p in folder.iterdir() if p.suffix.lower() in ALLOWED_EXTS]

    # Recurse if requested
    if recursive:
        for sub in folder.rglob("*"):
            if sub.is_dir() and sub != folder:
                process_folder(sub, recursive, dpi, quiet)

    # If both PDFs and images exist — prioritize PDFs, warn user
    if pdfs and imgs:
        print(f"⚠️  Mixed content in {folder.name}: prioritizing PDFs.")
        targets = pdfs
    elif pdfs:
        targets = pdfs
    elif imgs:
        targets = [folder]
    else:
        if not quiet:
            print(f"(empty) {folder}")
        return

    for item in targets:
        if item.suffix.lower() == ".pdf":
            out_name = item.stem + ".pptx"
            out_path = folder / out_name
            print(f"📄 Converting PDF → PPTX: {item.name} → {out_name}")
            pages = convert_pdf_to_images(item, dpi=dpi)
            w_in, h_in = pdf_first_page_size_inches(item)
            build_presentation(pages, out_path, w_in, h_in, "fit", show_progress=True)
        else:
            # Image folder
            imgs = list_images(item if item.is_dir() else folder)
            if not imgs:
                continue
            out_name = folder.name + ".pptx"
            out_path = folder / out_name
            print(f"🖼️  Building PPTX from {len(imgs)} images → {out_name}")

            # Detect aspect ratio from first image
            from PIL import Image

            with Image.open(imgs[0]) as im:
                w_in, h_in = (
                    im.width / 96,
                    im.height / 96,
                )  # assume 96 DPI if not embedded
                # normalize: ensure landscape orientation for convenience
                if w_in < h_in:
                    w_in, h_in = h_in, w_in

            build_presentation(imgs, out_path, w_in, h_in, "fit", show_progress=True)


# ===[ MAIN ENTRYPOINT ]============================================


def main():
    """Entry point for pptx-builder CLI — supports CLI or interactive use."""
    args = parse_cli_args()

    # Configure logging based on verbose flag
    if args.verbose:
        logging.basicConfig(level=logging.DEBUG, format="%(levelname)s: %(message)s")
    else:
        logging.basicConfig(level=logging.WARNING, format="%(levelname)s: %(message)s")

    # Validate --output usage
    if args.output and args.input and len(args.input) > 1:
        print("✗ Error: --output can only be used with a single input file.")
        return

    # Interactive fallback if no input flag provided
    if not args.input:
        print("\n=== PPTX Builder (Interactive Mode) ===\n")
        in_path = prompt_input_path()

        kind = detect_input_type(in_path)
        if kind == "unknown":
            print("✗ That path is not a PDF or an image folder. Exiting.")
            sys.exit(1)

        out_name = prompt_output_name(default_name="slides")
        if in_path.is_dir():
            output_path = (in_path / out_name).resolve()
        else:
            # Use --output if provided, otherwise use input name
            if args.output:
                out_name = args.output
                if not out_name.lower().endswith(".pptx"):
                    out_name = out_name + ".pptx"
                output_path = in_path.parent / out_name
            else:
                output_path = in_path.with_suffix(".pptx")

        width_in, height_in = prompt_slide_size()
        mode = prompt_fit_mode()  # Let user choose fit or fill

        print("\nSummary:")
        print(f"  Source: {in_path}")
        file_type = "PDF file" if kind == "pdf" else "Image folder"
        print(f"  Type  : {file_type}")
        print(f'  Slide size: {width_in:.2f}" x {height_in:.2f}"')
        placement = "Fit whole image (no crop)" if mode == "fit" else "Crop to fill (no whitespace)"
        print(f"  Placement : {placement}")
        print(f"  Output file: {output_path}\n")

        temp_dir = None
        try:
            if kind == "pdf":
                pages = convert_pdf_to_images(in_path, dpi=args.dpi)
                # Store temp dir for cleanup
                if pages:
                    temp_dir = pages[0].parent
                build_presentation(
                    images=pages,
                    output_path=output_path,
                    slide_width_in=width_in,
                    slide_height_in=height_in,
                    mode=mode,
                    show_progress=True,
                )
            else:
                images = list_images(in_path)
                if not images:
                    print("No images found in that folder. Exiting.")
                    sys.exit(1)
                build_presentation(
                    images=images,
                    output_path=output_path,
                    slide_width_in=width_in,
                    slide_height_in=height_in,
                    mode=mode,
                    show_progress=True,
                )
        except Exception as e:
            print(f"✗ Failed to create presentation: {e}")
            sys.exit(1)
        finally:
            # Clean up temporary PDF conversion files
            if temp_dir and temp_dir.exists():
                import shutil

                try:
                    shutil.rmtree(temp_dir)
                except Exception:
                    pass  # Best effort cleanup

        print(f"✅ Presentation saved to: {output_path}")
        return

    # Non-interactive CLI mode
    for path_str in args.input:
        path = Path(path_str).expanduser().resolve()

        if not path.exists():
            print(f"✗ Input not found: {path}")
            continue

        kind = detect_input_type(path)

        if kind == "pdf":
            print(f"📄 [CLI] Converting PDF → PPTX: {path.name}")
            temp_dir = None
            try:
                pages = convert_pdf_to_images(path, dpi=args.dpi)
                if pages:
                    temp_dir = pages[0].parent

                # Determine output name
                if args.output:
                    out_name = args.output
                    if not out_name.lower().endswith(".pptx"):
                        out_name = out_name + ".pptx"
                else:
                    out_name = path.stem + ".pptx"

                out_path = path.parent / out_name

                # 🔒 Overwrite protection
                w_in, h_in = pdf_first_page_size_inches(path)
                if confirm_overwrite(out_path, quiet=args.quiet, force=args.force):
                    build_presentation(
                        pages,
                        output_path=out_path,
                        slide_width_in=w_in,
                        slide_height_in=h_in,
                        mode="fit",
                        show_progress=not args.quiet,
                    )
                    if not args.quiet:
                        print(f'✅ Saved: {out_path} ({w_in:.2f}" × {h_in:.2f}")')
                else:
                    print(f"⏩ Skipped (already exists): {out_path}")

            except Exception as e:
                print(f"✗ Failed to process {path}: {e}")
            finally:
                # Clean up temp files
                if temp_dir and temp_dir.exists():
                    import shutil

                    try:
                        shutil.rmtree(temp_dir)
                    except Exception:
                        pass

        elif kind == "folder":
            if not args.quiet:
                print(f"🗂️  [CLI] Processing folder: {path}")
            try:
                process_folder(path, recursive=args.recursive, dpi=args.dpi, quiet=args.quiet)
            except Exception as e:
                print(f"✗ Folder failed: {path} ({e})")

        else:
            print(f"✗ Unsupported input: {path}")

    if not args.quiet:
        print("\n✅ CLI execution complete.")


if __name__ == "__main__":
    main()
