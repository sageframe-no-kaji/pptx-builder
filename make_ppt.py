
#!/usr/bin/env python3
# SIMPLE VERSION - IN PRODUCTION
"""
make_ppt.py
-----------
Create a PowerPoint (.pptx) from a folder of images (PNG/JPG/JPEG).

Features:
    - Prompts for source folder and output filename.
    - Lets you choose common slide sizes (16:9, 4:3, Letter, A4, Legal, Tabloid).
    - Two placement modes:
        (1) Fit whole image (no crop, may show background bars),
        (2) Crop-to-fill (cover slide fully, cropping as needed; no distortion).
    - One image per slide, centered, sorted by filename (case-insensitive).

Notes:
    - Uses python-pptx (and Pillow) under the hood.
    - Slide background shows through when using "Fit whole image".
    - No part of the image is stretched; scaling is always proportional.
"""

import sys
import os
from pathlib import Path
from typing import Tuple, List

from pptx import Presentation
from pptx.util import Inches, Emu


# -----------------------------
# Slide size presets (in inches)
# -----------------------------
SLIDE_SIZES = {
    "1": ("16:9 (13.33\" x 7.5\")", 13.3333333333, 7.5),    # PowerPoint default widescreen
    "2": ("4:3  (10\" x 7.5\")",     10.0, 7.5),
    "3": ("Letter  (11\" x 8.5\")",  11.0, 8.5),
    "4": ("A4      (11.69\" x 8.27\")", 11.69, 8.27),
    "5": ("Legal   (14\" x 8.5\")",  14.0, 8.5),
    "6": ("Tabloid (17\" x 11\")",   17.0, 11.0),
}

# -----------------------------
# File extensions we will accept
# -----------------------------
ALLOWED_EXTS = {".png", ".jpg", ".jpeg"}


def prompt_folder() -> Path:
    """Ask the user for the path to the folder containing images."""
    while True:
        folder_str = input("Enter the path to the folder containing images (PNG/JPG): ").strip()
        folder = Path(folder_str).expanduser().resolve()
        if folder.is_dir():
            return folder
        print("✗ That path is not a folder. Please try again.\n")


def prompt_output_name(default_name: str = "output") -> Path:
    """Ask the user for the output file name (no extension needed)."""
    out = input(f"Enter output filename (without extension) [{default_name}]: ").strip()
    if not out:
        out = default_name
    # Ensure .pptx extension
    if not out.lower().endswith(".pptx"):
        out = out + ".pptx"
    return Path(out).name  # return just the name (no path), we'll place it next to the script or in the folder later


def prompt_slide_size() -> Tuple[float, float]:
    """Display size options and return (width_in_inches, height_in_inches)."""
    print("\nChoose slide size:")
    for key, (label, w, h) in SLIDE_SIZES.items():
        print(f"  {key}) {label}")
    while True:
        choice = input("Enter number (1-6): ").strip()
        if choice in SLIDE_SIZES:
            _, w, h = SLIDE_SIZES[choice]
            return w, h
        print("✗ Invalid choice. Please enter a number from the list.\n")


def prompt_fit_mode() -> str:
    """
    Ask how images should be placed:
        1) Fit whole image (no cropping)  -> 'fit'
        2) Crop to fill (cover, may crop) -> 'fill'
    """
    print("\nHow should images be placed?")
    print("  1) Fit whole image (no cropping; background may show)")
    print("  2) Crop to fill (no whitespace; edges may be trimmed)")
    while True:
        choice = input("Enter 1 or 2: ").strip()
        if choice == "1":
            return "fit"
        if choice == "2":
            return "fill"
        print("✗ Invalid choice. Please enter 1 or 2.\n")


def list_images(folder: Path) -> List[Path]:
    """Return sorted list of image files with allowed extensions (case-insensitive)."""
    files = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in ALLOWED_EXTS]
    # Sort case-insensitively by filename
    files.sort(key=lambda p: p.name.lower())
    return files


def emu_to_float_inches(emu: Emu) -> float:
    """Convert EMU to inches (pptx.util.Inches wraps conversion, but we need a float)."""
    # 1 inch = 914400 EMU
    return float(emu) / 914400.0


def place_picture_fit(slide, img_path: Path, slide_w_emu: int, slide_h_emu: int):
    """
    Place the image on the slide using 'contain' behavior:
        - Scale proportionally so the entire image is visible (no cropping).
        - Center the image; background may show (letterbox/pillarbox).
    Implementation detail:
        - Insert at natural size first to read pic.width/height (requires Pillow via python-pptx).
        - Compute scale ratio and then set final size and position.
    """
    pic = slide.shapes.add_picture(str(img_path), left=0, top=0)  # natural size first
    img_w = float(pic.width)
    img_h = float(pic.height)
    sw = float(slide_w_emu)
    sh = float(slide_h_emu)

    # Scale to "contain": use the smaller ratio
    scale = min(sw / img_w, sh / img_h)
    new_w = img_w * scale
    new_h = img_h * scale

    # Center
    left = (sw - new_w) / 2.0
    top = (sh - new_h) / 2.0

    pic.left = int(left)
    pic.top = int(top)
    pic.width = int(new_w)
    pic.height = int(new_h)


def place_picture_fill(slide, img_path: Path, slide_w_emu: int, slide_h_emu: int):
    """
    Place the image on the slide using 'cover' behavior:
        - Scale proportionally so the slide is fully covered.
        - May crop the image (overflow outside slide bounds is not visible).
        - Center the image.
    """
    pic = slide.shapes.add_picture(str(img_path), left=0, top=0)  # natural size first
    img_w = float(pic.width)
    img_h = float(pic.height)
    sw = float(slide_w_emu)
    sh = float(slide_h_emu)

    # Scale to "cover": use the larger ratio
    scale = max(sw / img_w, sh / img_h)
    new_w = img_w * scale
    new_h = img_h * scale

    # Center (image may overflow the slide; that's fine)
    left = (sw - new_w) / 2.0
    top = (sh - new_h) / 2.0

    pic.left = int(left)
    pic.top = int(top)
    pic.width = int(new_w)
    pic.height = int(new_h)


def build_presentation(
    images: List[Path],
    output_path: Path,
    slide_width_in: float,
    slide_height_in: float,
    mode: str
) -> None:
    """Create the PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    # Cache slide size in EMU for math
    sw_emu = int(prs.slide_width)
    sh_emu = int(prs.slide_height)

    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        if mode == "fit":
            place_picture_fit(slide, img, sw_emu, sh_emu)
        else:
            place_picture_fill(slide, img, sw_emu, sh_emu)

    prs.save(str(output_path))


# ===[ SECTION: CLI ARGUMENTS ]=====================================

import argparse

def parse_cli_args():
    """Parse command-line arguments for batch, recursive, or quiet runs."""
    parser = argparse.ArgumentParser(
        prog="make_ppt.py",
        description="Build PowerPoint (.pptx) files from PDFs or image folders."
    )

    parser.add_argument(
        "-i", "--input", nargs="+", metavar="PATH",
        help="Path(s) to one or more PDFs, images, or folders to process."
    )

    parser.add_argument(
        "-r", "--recursive", action="store_true",
        help="Recurse into subfolders when processing folders."
    )

    parser.add_argument(
        "--dpi", type=int, default=300,
        help="DPI for PDF rendering (default: 300)."
    )

    parser.add_argument(
        "--quiet", action="store_true",
        help="Suppress interactive prompts and non-critical output."
    )

    return parser.parse_args()

# ===[ SECTION: INPUT HANDLING ]====================================

from pdf2image import convert_from_path


def detect_input_type(path: Path) -> str:
    """Return 'pdf', 'folder', or 'unknown' based on the given path."""
    if path.is_file() and path.suffix.lower() == ".pdf":
        return "pdf"
    if path.is_dir():
        # check if folder contains images
        imgs = list_images(path)
        if imgs:
            return "folder"
    return "unknown"


def convert_pdf_to_images(pdf_path: Path, dpi: int) -> List[Path]:
    """Convert PDF pages to temporary PNG files at given DPI; return list of paths."""
    import tempfile
    from PIL import Image

    temp_dir = Path(tempfile.mkdtemp(prefix="pptx_pdf_"))
    pages = convert_from_path(pdf_path.as_posix(), dpi=dpi)
    out_paths = []

    for i, page in enumerate(pages, start=1):
        out_path = temp_dir / f"page_{i:04d}.png"
        page.save(out_path, "PNG")
        out_paths.append(out_path)

    return out_paths


def process_folder(folder: Path, recursive: bool, dpi: int, quiet: bool) -> None:
    """Process all PDFs and/or images in a folder into PPTX files."""
    pdfs = sorted(folder.glob("*.pdf"))
    imgs = [p for p in folder.iterdir() if p.suffix.lower() in ALLOWED_EXTS]

    # Recurse if requested
    if recursive:
        for sub in folder.rglob("*"):
            if sub.is_dir() and sub != folder:
                process_folder(sub, recursive, dpi, quiet)

    # If both PDFs and images exist — prioritize PDFs, warn user
    if pdfs and imgs:
        print(f"⚠️  Mixed content in {folder.name}: prioritizing PDFs.")
        targets = pdfs
    elif pdfs:
        targets = pdfs
    elif imgs:
        targets = [folder]
    else:
        if not quiet:
            print(f"(empty) {folder}")
        return

    for item in targets:
        if item.suffix.lower() == ".pdf":
            out_name = item.stem + ".pptx"
            out_path = folder / out_name
            print(f"📄 Converting PDF → PPTX: {item.name} → {out_name}")
            pages = convert_pdf_to_images(item, dpi=dpi)
            build_presentation(pages, out_path, slide_size_in=(13.3333, 7.5), mode="fit")
        else:
            # Image folder
            imgs = list_images(item if item.is_dir() else folder)
            if not imgs:
                continue
            out_name = folder.name + ".pptx"
            out_path = folder / out_name
            print(f"🖼️  Building PPTX from {len(imgs)} images → {out_name}")
            build_presentation(imgs, out_path, slide_size_in=(13.3333, 7.5), mode="fit")

# ===[ MAIN ENTRYPOINT ]============================================

def main():
    """Entry point for make_ppt.py — supports CLI or interactive use."""
    args = parse_cli_args()

    # Interactive fallback if no input flag provided
    if not args.input:
        print("\n=== PPTX from Images ===\n")
        folder = prompt_folder()
        images = list_images(folder)
        if not images:
            print("No PNG/JPG images found in that folder. Exiting.")
            sys.exit(1)

        out_name = prompt_output_name(default_name="slides")
        output_path = (folder / out_name).resolve()

        width_in, height_in = prompt_slide_size()
        mode = "fit"  # always fit mode by design

        print("\nSummary:")
        print(f"  Source folder: {folder}")
        print(f"  Images found : {len(images)}")
        print(f"  Slide size   : {width_in:.2f}\" x {height_in:.2f}\"")
        print(f"  Placement    : Fit whole image (no crop)")
        print(f"  Output file  : {output_path}\n")

        try:
            build_presentation(
                images=images,
                output_path=output_path,
                slide_width_in=width_in,
                slide_height_in=height_in,
                mode=mode
            )
        except Exception as e:
            print(f"✗ Failed to create presentation: {e}")
            sys.exit(1)

        print(f"✅ Presentation saved to: {output_path}")
        return

    # Non-interactive CLI mode
    for path_str in args.input:
        path = Path(path_str).expanduser().resolve()

        if not path.exists():
            print(f"✗ Input not found: {path}")
            continue

        kind = detect_input_type(path)

        if kind == "pdf":
            print(f"📄 [CLI] Converting PDF → PPTX: {path.name}")
            try:
                pages = convert_pdf_to_images(path, dpi=args.dpi)
                out_name = path.stem + ".pptx"
                out_path = path.parent / out_name
                build_presentation(
                    pages,
                    output_path=out_path,
                    slide_width_in=13.3333,
                    slide_height_in=7.5,
                    mode="fit"
                )
                if not args.quiet:
                    print(f"✅ Saved: {out_path}")
            except Exception as e:
                print(f"✗ Failed to process {path}: {e}")

        elif kind == "folder":
            if not args.quiet:
                print(f"🗂️  [CLI] Processing folder: {path}")
            try:
                process_folder(path, recursive=args.recursive, dpi=args.dpi, quiet=args.quiet)
            except Exception as e:
                print(f"✗ Folder failed: {path} ({e})")

        else:
            print(f"✗ Unsupported input: {path}")

    if not args.quiet:
        print("\n✅ CLI execution complete.")


if __name__ == "__main__":
    main()
