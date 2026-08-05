"""
Tests for pptx_builder.core and pptx_builder.web
"""

import sys
import time
import shutil
import tempfile
import logging
from contextlib import ExitStack
from pathlib import Path
from unittest.mock import patch, MagicMock, PropertyMock
from PIL import Image

import pytest
from pptx.util import Emu, Inches

from pptx_builder.core import (
    list_images,
    detect_input_type,
    confirm_overwrite,
    emu_to_float_inches,
    build_presentation,
    place_picture_fit,
    place_picture_fill,
    prompt_input_path,
    prompt_output_name,
    prompt_slide_size,
    prompt_fit_mode,
    convert_pdf_to_images,
    pdf_first_page_size_inches,
    process_folder,
    parse_cli_args,
    prompt_output_format,
    main,
    ALLOWED_EXTS,
    SLIDE_SIZES,
    OUTPUT_FORMATS,
    DEFAULT_DPI,
    DEFAULT_FORMAT,
    DEFAULT_QUALITY,
    _normalize_format,
    _env_int,
    _env_format,
)
import pptx_builder.web as web_module
from pptx_builder.web import (
    process_files,
    cleanup_temp_files,
    cleanup_old_files,
    _build_info_strip,
    _quality_visibility,
    SLIDE_SIZE_OPTIONS,
    MAX_FILE_SIZE,
    MAX_FILES,
    MAX_QUALITY,
)
import gradio as gr

# ─── helpers ─────────────────────────────────────────────────────────────────


def _make_image(path: Path, width: int = 100, height: int = 100, color: str = "red") -> Path:
    Image.new("RGB", (width, height), color=color).save(path)
    return path


def _make_pdf(path: Path, pages: int = 3, width: int = 612, height: int = 792) -> Path:
    """Write a real multi-page PDF carrying photographic content via PyMuPDF.

    The page content is seeded pseudo-random noise, embedded as a raster image.
    Content shape matters here: flat fills and hard-edged bands compress better
    under PNG's row filters than under JPEG's DCT, so a synthetic banded page
    would show JPEG *larger* and misrepresent what real slides do. Noise stands
    in for the photographic content that makes JPEG the smaller encoder.
    """
    import math
    import random

    import fitz

    rng = random.Random(1729)
    # Source raster is sized to the rendered page so it is not upscaled — an
    # upscaled image renders as hard-edged blocks, which PNG compresses well and
    # JPEG does not, inverting the property under test.
    px_w, px_h = int(width * 1.4), int(height * 1.4)
    doc = fitz.open()
    for n in range(pages):
        # Smooth two-dimensional gradient plus fine jitter: the texture profile of
        # a photograph. Smoothness is what JPEG's DCT compresses well, and the
        # per-pixel jitter is what defeats PNG's row filters.
        buf = bytearray()
        for y in range(px_h):
            fy = y / px_h
            for x in range(px_w):
                fx = x / px_w
                base = 128 + 100 * math.sin(fx * 3.1 + n) * math.cos(fy * 2.7)
                for shift in (0.0, 2.1, 4.2):
                    value = base + 40 * math.sin(fx * 5.0 + shift) + rng.uniform(-6, 6)
                    buf.append(max(0, min(255, int(value))))
        photo = Image.frombytes("RGB", (px_w, px_h), bytes(buf))
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as handle:
            photo_path = Path(handle.name)
        photo.save(photo_path)
        try:
            page = doc.new_page(width=width, height=height)
            page.insert_image(fitz.Rect(0, 0, width, height), filename=str(photo_path))
        finally:
            photo_path.unlink()
    doc.save(str(path))
    doc.close()
    return path


def _make_fitz_mock(width: float = 612.0, height: float = 792.0, page_count: int = 1):
    mock_fitz = MagicMock()
    mock_doc = MagicMock()
    mock_doc.page_count = page_count
    mock_page = MagicMock()
    mock_page.rect.width = width
    mock_page.rect.height = height
    mock_doc.__getitem__ = MagicMock(return_value=mock_page)
    mock_doc.__enter__ = MagicMock(return_value=mock_doc)
    mock_doc.__exit__ = MagicMock(return_value=False)
    mock_fitz.open.return_value = mock_doc
    return mock_fitz, mock_doc


@pytest.fixture(autouse=True)
def _clean_temp_dirs():
    original = list(web_module.TEMP_DIRS)
    web_module.TEMP_DIRS.clear()
    yield
    web_module.TEMP_DIRS.clear()
    web_module.TEMP_DIRS.extend(original)


# ─── image listing ────────────────────────────────────────────────────────────


class TestImageListing:
    def test_filters_correctly(self, tmp_path):
        (tmp_path / "image.png").touch()
        (tmp_path / "photo.jpg").touch()
        (tmp_path / "doc.txt").touch()
        images = list_images(tmp_path)
        assert len(images) == 2
        assert all(p.suffix.lower() in ALLOWED_EXTS for p in images)

    def test_case_insensitive(self, tmp_path):
        (tmp_path / "IMAGE.PNG").touch()
        (tmp_path / "PHOTO.JPG").touch()
        assert len(list_images(tmp_path)) == 2

    def test_sorted_alphabetically(self, tmp_path):
        for name in ("zebra.png", "Apple.jpg", "banana.png"):
            (tmp_path / name).touch()
        names = [p.name.lower() for p in list_images(tmp_path)]
        assert names == sorted(names)


# ─── input detection ──────────────────────────────────────────────────────────


class TestInputDetection:
    def test_detect_pdf(self, tmp_path):
        f = tmp_path / "t.pdf"
        f.touch()
        assert detect_input_type(f) == "pdf"

    def test_detect_folder_with_images(self, tmp_path):
        (tmp_path / "img.png").touch()
        assert detect_input_type(tmp_path) == "folder"

    def test_detect_empty_folder(self, tmp_path):
        assert detect_input_type(tmp_path) == "unknown"

    def test_detect_unknown_file(self, tmp_path):
        f = tmp_path / "t.txt"
        f.touch()
        assert detect_input_type(f) == "unknown"


# ─── slide presets ────────────────────────────────────────────────────────────


class TestSlidePresets:
    def test_all_presets_valid(self):
        for _, (label, w, h) in SLIDE_SIZES.items():
            assert w > 0 and h > 0 and len(label) > 0

    def test_standard_sizes_included(self):
        labels = [l for l, _, _ in SLIDE_SIZES.values()]
        assert any("16:9" in l for l in labels)
        assert any("4:3" in l for l in labels)
        assert any("Letter" in l for l in labels)
        assert any("A4" in l for l in labels)


# ─── allowed extensions ───────────────────────────────────────────────────────


class TestAllowedExtensions:
    def test_common_formats_supported(self):
        assert {".png", ".jpg", ".jpeg"}.issubset(ALLOWED_EXTS)

    def test_heic_supported(self):
        assert ".heic" in ALLOWED_EXTS or ".heif" in ALLOWED_EXTS

    def test_extensions_lowercase(self):
        assert all(e.islower() for e in ALLOWED_EXTS)


# ─── utilities ────────────────────────────────────────────────────────────────


class TestUtilities:
    def test_emu_conversion(self):
        assert abs(emu_to_float_inches(Emu(914400)) - 1.0) < 0.001
        assert abs(emu_to_float_inches(Emu(1828800)) - 2.0) < 0.001
        assert emu_to_float_inches(Emu(0)) == 0.0

    def test_confirm_overwrite_force(self, tmp_path):
        f = tmp_path / "e.pptx"
        f.touch()
        assert confirm_overwrite(f, force=True) is True

    def test_confirm_overwrite_quiet(self, tmp_path):
        f = tmp_path / "e.pptx"
        f.touch()
        assert confirm_overwrite(f, quiet=True) is True

    def test_confirm_overwrite_nonexistent(self, tmp_path):
        nf = tmp_path / "new.pptx"
        assert confirm_overwrite(nf, quiet=True) is True
        assert confirm_overwrite(nf, force=True) is True

    def test_confirm_overwrite_interactive_yes(self, tmp_path):
        f = tmp_path / "e.pptx"
        f.touch()
        with patch("builtins.input", return_value="y"):
            assert confirm_overwrite(f) is True

    def test_confirm_overwrite_interactive_no(self, tmp_path):
        f = tmp_path / "e.pptx"
        f.touch()
        with patch("builtins.input", return_value="n"):
            assert confirm_overwrite(f) is False


# ─── build_presentation ───────────────────────────────────────────────────────


class TestBuildPresentation:
    @pytest.mark.integration
    def test_fit_mode(self, tmp_path):
        img = _make_image(tmp_path / "t.png", 200, 100)
        out = tmp_path / "fit.pptx"
        build_presentation([img], out, 10.0, 7.5, "fit")
        assert out.stat().st_size > 10000

    @pytest.mark.integration
    def test_fill_mode(self, tmp_path):
        img = _make_image(tmp_path / "t.png", 300, 200)
        out = tmp_path / "fill.pptx"
        build_presentation([img], out, 10.0, 7.5, "fill")
        assert out.stat().st_size > 10000

    @pytest.mark.integration
    def test_multiple_images(self, tmp_path):
        imgs = [
            _make_image(tmp_path / f"i{n}.png", color=c)
            for n, c in enumerate(["red", "green", "blue"])
        ]
        out = tmp_path / "multi.pptx"
        build_presentation(imgs, out, 10.0, 7.5, "fit")
        assert out.exists()

    @pytest.mark.integration
    def test_custom_dimensions(self, tmp_path):
        img = _make_image(tmp_path / "t.png")
        out = tmp_path / "a4.pptx"
        build_presentation([img], out, 11.69, 8.27, "fit")
        assert out.exists()

    @pytest.mark.integration
    def test_show_progress_true(self, tmp_path, capsys):
        img = _make_image(tmp_path / "t.png")
        out = tmp_path / "prog.pptx"
        build_presentation([img], out, 10.0, 7.5, "fit", show_progress=True)
        assert out.exists()
        assert "Saving" in capsys.readouterr().out


# ─── place_picture zero-dimension guards ─────────────────────────────────────


class TestPlacePictureGuards:
    def _slide(self):
        from pptx import Presentation

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _mock_pic(self, w, h):
        p = MagicMock()
        p.width = w
        p.height = h
        return p

    def test_fit_zero_width(self, tmp_path):
        slide = self._slide()
        img = tmp_path / "i.png"
        img.touch()
        with patch.object(slide.shapes, "add_picture", return_value=self._mock_pic(0, 100)):
            with pytest.raises(ValueError, match="zero dimensions"):
                place_picture_fit(slide, img, 9144000, 6858000)

    def test_fit_zero_height(self, tmp_path):
        slide = self._slide()
        img = tmp_path / "i.png"
        img.touch()
        with patch.object(slide.shapes, "add_picture", return_value=self._mock_pic(100, 0)):
            with pytest.raises(ValueError, match="zero dimensions"):
                place_picture_fit(slide, img, 9144000, 6858000)

    def test_fill_zero_width(self, tmp_path):
        slide = self._slide()
        img = tmp_path / "i.png"
        img.touch()
        with patch.object(slide.shapes, "add_picture", return_value=self._mock_pic(0, 100)):
            with pytest.raises(ValueError, match="zero dimensions"):
                place_picture_fill(slide, img, 9144000, 6858000)

    def test_fill_zero_height(self, tmp_path):
        slide = self._slide()
        img = tmp_path / "i.png"
        img.touch()
        with patch.object(slide.shapes, "add_picture", return_value=self._mock_pic(100, 0)):
            with pytest.raises(ValueError, match="zero dimensions"):
                place_picture_fill(slide, img, 9144000, 6858000)


# ─── prompt functions ─────────────────────────────────────────────────────────


class TestPromptFunctions:
    def test_input_path_valid(self, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        with patch("builtins.input", return_value=str(pdf)):
            assert prompt_input_path() == pdf.resolve()

    def test_input_path_retries(self, tmp_path):
        pdf = tmp_path / "r.pdf"
        pdf.touch()
        with patch("builtins.input", side_effect=["/nonexistent/x.pdf", str(pdf)]):
            assert prompt_input_path() == pdf.resolve()

    def test_output_name_with_extension(self):
        with patch("builtins.input", return_value="slides.pptx"):
            assert prompt_output_name() == "slides.pptx"

    def test_output_name_without_extension(self):
        with patch("builtins.input", return_value="my_slides"):
            assert prompt_output_name() == "my_slides.pptx"

    def test_output_name_empty_uses_default(self):
        with patch("builtins.input", return_value=""):
            assert prompt_output_name(default_name="out") == "out.pptx"

    def test_slide_size_choice_1(self):
        with patch("builtins.input", return_value="1"):
            w, h = prompt_slide_size()
        assert abs(w - 13.3333) < 0.01 and h == 7.5

    def test_slide_size_retries_invalid(self):
        with patch("builtins.input", side_effect=["9", "2"]):
            w, h = prompt_slide_size()
        assert w == 10.0 and h == 7.5

    def test_fit_mode_returns_fit(self):
        with patch("builtins.input", return_value="1"):
            assert prompt_fit_mode() == "fit"

    def test_fit_mode_returns_fill(self):
        with patch("builtins.input", return_value="2"):
            assert prompt_fit_mode() == "fill"

    def test_fit_mode_retries_invalid(self):
        with patch("builtins.input", side_effect=["x", "1"]):
            assert prompt_fit_mode() == "fit"


# ─── parse_cli_args ───────────────────────────────────────────────────────────


class TestParseCLIArgs:
    def test_defaults(self):
        with patch("sys.argv", ["pb"]):
            args = parse_cli_args()
        assert args.input is None
        assert args.dpi == 200
        assert args.format == "jpeg"
        assert args.quality == 85
        assert not args.quiet
        assert not args.verbose
        assert not args.force
        assert args.output is None

    def test_input_flag(self):
        with patch("sys.argv", ["pb", "-i", "a.pdf", "b.pdf"]):
            assert parse_cli_args().input == ["a.pdf", "b.pdf"]

    def test_dpi_flag(self):
        with patch("sys.argv", ["pb", "--dpi", "150"]):
            assert parse_cli_args().dpi == 150

    def test_quiet_flag(self):
        with patch("sys.argv", ["pb", "--quiet"]):
            assert parse_cli_args().quiet is True

    def test_verbose_flag(self):
        with patch("sys.argv", ["pb", "--verbose"]):
            assert parse_cli_args().verbose is True

    def test_force_flag(self):
        with patch("sys.argv", ["pb", "--force"]):
            assert parse_cli_args().force is True

    def test_output_flag(self):
        with patch("sys.argv", ["pb", "-o", "out.pptx"]):
            assert parse_cli_args().output == "out.pptx"

    def test_recursive_flag(self):
        with patch("sys.argv", ["pb", "-r"]):
            assert parse_cli_args().recursive is True


# ─── convert_pdf_to_images ────────────────────────────────────────────────────


class TestConvertPDFToImages:
    @patch("pptx_builder.core.convert_from_path")
    def test_success_returns_paths(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()

        def fake(path, dpi, output_folder, fmt, paths_only, **kwargs):
            p = Path(output_folder) / f"page-0001.{fmt}"
            p.touch()
            return [str(p)]

        mock_cv.side_effect = fake
        result = convert_pdf_to_images(pdf, dpi=150)
        assert len(result) == 1
        assert isinstance(result[0], Path)

    @patch("pptx_builder.core.convert_from_path")
    def test_failure_cleans_dir_and_raises(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.side_effect = Exception("poppler missing")
        with pytest.raises(RuntimeError, match="Failed to convert PDF"):
            convert_pdf_to_images(pdf, dpi=150)

    @patch("pptx_builder.core.convert_from_path")
    def test_uses_paths_only_flag(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.return_value = []
        convert_pdf_to_images(pdf, dpi=300)
        _, kw = mock_cv.call_args
        assert kw.get("paths_only") is True

    @patch("pptx_builder.core.convert_from_path")
    def test_empty_pdf_returns_empty_list(self, mock_cv, tmp_path):
        pdf = tmp_path / "e.pdf"
        pdf.touch()
        mock_cv.return_value = []
        assert convert_pdf_to_images(pdf, dpi=150) == []

    @patch("pptx_builder.core.convert_from_path")
    def test_max_pages_passes_last_page(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.return_value = []
        convert_pdf_to_images(pdf, dpi=150, max_pages=5)
        _, kw = mock_cv.call_args
        assert kw.get("last_page") == 5

    @patch("pptx_builder.core.convert_from_path")
    def test_max_pages_zero_omits_last_page(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.return_value = []
        convert_pdf_to_images(pdf, dpi=150, max_pages=0)
        _, kw = mock_cv.call_args
        assert "last_page" not in kw


# ─── _build_info_strip ────────────────────────────────────────────────────────


class TestOutputFormatConfig:
    """Environment/flag resolution for the output encoder (ho-01)."""

    def test_normalize_folds_jpg_alias(self):
        assert _normalize_format("JPG") == "jpeg"
        assert _normalize_format(" Jpeg ") == "jpeg"
        assert _normalize_format("PNG") == "png"

    def test_builtin_defaults(self):
        # Guards the ho-01 decision: bare invocation is 200 DPI JPEG q85.
        assert DEFAULT_DPI == 200
        assert DEFAULT_FORMAT == "jpeg"
        assert DEFAULT_QUALITY == 85

    def test_env_int_reads_and_falls_back(self, monkeypatch):
        monkeypatch.setenv("PPTX_TEST_INT", "42")
        assert _env_int("PPTX_TEST_INT", 7, minimum=1, maximum=100) == 42
        # Malformed and out-of-range both fall back rather than raising.
        monkeypatch.setenv("PPTX_TEST_INT", "banana")
        assert _env_int("PPTX_TEST_INT", 7, minimum=1, maximum=100) == 7
        monkeypatch.setenv("PPTX_TEST_INT", "999")
        assert _env_int("PPTX_TEST_INT", 7, minimum=1, maximum=100) == 7

    def test_env_format_reads_and_falls_back(self, monkeypatch):
        monkeypatch.setenv("PPTX_TEST_FMT", "png")
        assert _env_format("PPTX_TEST_FMT", "jpeg") == "png"
        monkeypatch.setenv("PPTX_TEST_FMT", "jpg")
        assert _env_format("PPTX_TEST_FMT", "png") == "jpeg"
        monkeypatch.setenv("PPTX_TEST_FMT", "tiff")
        assert _env_format("PPTX_TEST_FMT", "jpeg") == "jpeg"

    def test_env_absent_uses_fallback(self, monkeypatch):
        monkeypatch.delenv("PPTX_TEST_ABSENT", raising=False)
        assert _env_int("PPTX_TEST_ABSENT", 5, minimum=1, maximum=10) == 5
        assert _env_format("PPTX_TEST_ABSENT", "png") == "png"

    def test_flag_overrides_environment(self, monkeypatch):
        # Precedence: explicit flag beats environment beats built-in default.
        monkeypatch.setenv("PPTX_FORMAT", "png")
        monkeypatch.setenv("PPTX_QUALITY", "60")
        import importlib

        import pptx_builder.core as core_mod

        importlib.reload(core_mod)
        try:
            assert core_mod.DEFAULT_FORMAT == "png"
            assert core_mod.DEFAULT_QUALITY == 60
            with patch("sys.argv", ["pb"]):
                assert core_mod.parse_cli_args().format == "png"
            with patch("sys.argv", ["pb", "--format", "jpeg", "--quality", "95"]):
                args = core_mod.parse_cli_args()
            assert args.format == "jpeg"
            assert args.quality == 95
        finally:
            monkeypatch.delenv("PPTX_FORMAT", raising=False)
            monkeypatch.delenv("PPTX_QUALITY", raising=False)
            importlib.reload(core_mod)

    def test_quality_out_of_range_exits(self):
        for bad in ("0", "101", "abc"):
            with patch("sys.argv", ["pb", "--quality", bad]):
                with pytest.raises(SystemExit) as exc:
                    parse_cli_args()
            assert exc.value.code != 0

    def test_quality_accepted_with_png(self):
        # Accepted and ignored, so scripted callers can pass both unconditionally.
        with patch("sys.argv", ["pb", "--format", "png", "--quality", "50"]):
            args = parse_cli_args()
        assert args.format == "png"
        assert args.quality == 50

    def test_format_choices_rejects_unknown(self):
        with patch("sys.argv", ["pb", "--format", "tiff"]):
            with pytest.raises(SystemExit):
                parse_cli_args()

    def test_jpg_alias_accepted_on_cli(self):
        with patch("sys.argv", ["pb", "--format", "jpg"]):
            assert parse_cli_args().format == "jpeg"


class TestOutputFormatEncoding:
    """The encoder actually reaches Poppler and changes the bytes on disk."""

    @patch("pptx_builder.core.convert_from_path")
    def test_jpegopt_passed_for_jpeg(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.return_value = []
        convert_pdf_to_images(pdf, dpi=100, fmt="jpeg", quality=70)
        kwargs = mock_cv.call_args.kwargs
        assert kwargs["fmt"] == "jpeg"
        assert kwargs["jpegopt"] == {"quality": 70, "optimize": True, "progressive": True}

    @patch("pptx_builder.core.convert_from_path")
    def test_jpegopt_absent_for_png(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        mock_cv.return_value = []
        convert_pdf_to_images(pdf, dpi=100, fmt="png", quality=70)
        kwargs = mock_cv.call_args.kwargs
        assert kwargs["fmt"] == "png"
        assert "jpegopt" not in kwargs

    @patch("pptx_builder.core.convert_from_path")
    def test_unsupported_format_raises(self, mock_cv, tmp_path):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        with pytest.raises(ValueError, match="Unsupported output format"):
            convert_pdf_to_images(pdf, dpi=100, fmt="tiff")
        mock_cv.assert_not_called()

    @pytest.mark.integration
    def test_real_encoding_magic_bytes(self, tmp_path):
        # Assert on actual bytes, not the filename suffix.
        pdf = _make_pdf(tmp_path / "real.pdf", pages=2)
        jpeg_pages = convert_pdf_to_images(pdf, dpi=72, fmt="jpeg", quality=85)
        png_pages = convert_pdf_to_images(pdf, dpi=72, fmt="png")
        try:
            assert jpeg_pages and png_pages
            assert jpeg_pages[0].read_bytes()[:3] == b"\xff\xd8\xff"
            assert png_pages[0].read_bytes()[:4] == b"\x89PNG"
        finally:
            shutil.rmtree(jpeg_pages[0].parent, ignore_errors=True)
            shutil.rmtree(png_pages[0].parent, ignore_errors=True)

    @pytest.mark.integration
    def test_jpeg_smaller_than_png_at_matched_dpi(self, tmp_path):
        # Relative comparison only — absolute byte counts depend on Poppler build.
        pdf = _make_pdf(tmp_path / "sized.pdf", pages=2)
        jpeg_pages = convert_pdf_to_images(pdf, dpi=100, fmt="jpeg", quality=85)
        png_pages = convert_pdf_to_images(pdf, dpi=100, fmt="png")
        try:
            jpeg_bytes = sum(p.stat().st_size for p in jpeg_pages)
            png_bytes = sum(p.stat().st_size for p in png_pages)
            assert jpeg_bytes < png_bytes
        finally:
            shutil.rmtree(jpeg_pages[0].parent, ignore_errors=True)
            shutil.rmtree(png_pages[0].parent, ignore_errors=True)


class TestPromptOutputFormat:
    def test_jpeg_with_default_quality(self):
        with patch("builtins.input", side_effect=["1", ""]):
            assert prompt_output_format() == ("jpeg", DEFAULT_QUALITY)

    def test_empty_choice_defaults_to_jpeg(self):
        with patch("builtins.input", side_effect=["", ""]):
            assert prompt_output_format() == ("jpeg", DEFAULT_QUALITY)

    def test_png_skips_quality_prompt(self):
        # Only one input consumed — no quality question for PNG.
        with patch("builtins.input", side_effect=["2"]):
            assert prompt_output_format() == ("png", DEFAULT_QUALITY)

    def test_custom_quality(self):
        with patch("builtins.input", side_effect=["1", "60"]):
            assert prompt_output_format() == ("jpeg", 60)

    def test_reprompts_on_invalid(self, capsys):
        with patch("builtins.input", side_effect=["9", "1", "abc", "500", "75"]):
            assert prompt_output_format() == ("jpeg", 75)
        assert "Invalid choice" in capsys.readouterr().out

    def test_output_formats_constant(self):
        assert OUTPUT_FORMATS == ("png", "jpeg")


class TestBuildInfoStrip:
    def test_shows_configured_mb_and_files(self):
        import pptx_builder.web as wm

        orig_size, orig_files = wm.MAX_FILE_SIZE, wm.MAX_FILES
        orig_dpi, orig_pages = wm.MAX_DPI, wm.MAX_PDF_PAGES
        try:
            wm.MAX_FILE_SIZE = 20 * 1024 * 1024
            wm.MAX_FILES = 20
            wm.MAX_DPI = 300
            wm.MAX_PDF_PAGES = 20
            html = _build_info_strip()
            assert "20&nbsp;MB" in html
            assert "20 files" in html
            assert "300&nbsp;DPI" in html
            assert "20 pages" in html
        finally:
            wm.MAX_FILE_SIZE = orig_size
            wm.MAX_FILES = orig_files
            wm.MAX_DPI = orig_dpi
            wm.MAX_PDF_PAGES = orig_pages

    def test_omits_dpi_cap_when_at_max(self):
        import pptx_builder.web as wm

        orig_dpi = wm.MAX_DPI
        try:
            wm.MAX_DPI = 600
            html = _build_info_strip()
            assert "DPI max" not in html
        finally:
            wm.MAX_DPI = orig_dpi

    def test_omits_page_cap_when_unlimited(self):
        import pptx_builder.web as wm

        orig_pages = wm.MAX_PDF_PAGES
        try:
            wm.MAX_PDF_PAGES = 0
            html = _build_info_strip()
            assert "pages per PDF" not in html
        finally:
            wm.MAX_PDF_PAGES = orig_pages


# ─── pdf_first_page_size_inches ───────────────────────────────────────────────


class TestPDFFirstPageSize:
    def test_normal_pdf(self, tmp_path):
        mock_fitz, _ = _make_fitz_mock(612.0, 792.0, 1)
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            w, h = pdf_first_page_size_inches(pdf)
        assert abs(w - 8.5) < 0.01
        assert abs(h - 11.0) < 0.01

    def test_empty_pdf_returns_fallback(self, tmp_path):
        mock_fitz, _ = _make_fitz_mock(page_count=0)
        pdf = tmp_path / "e.pdf"
        pdf.touch()
        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            w, h = pdf_first_page_size_inches(pdf)
        assert abs(w - 13.3333) < 0.01
        assert h == 7.5

    def test_landscape_pdf(self, tmp_path):
        mock_fitz, _ = _make_fitz_mock(792.0, 612.0, 1)
        pdf = tmp_path / "l.pdf"
        pdf.touch()
        with patch.dict("sys.modules", {"fitz": mock_fitz}):
            w, h = pdf_first_page_size_inches(pdf)
        assert w > h


# ─── process_folder ───────────────────────────────────────────────────────────


class TestProcessFolder:
    @patch("pptx_builder.core.build_presentation")
    @patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
    @patch("pptx_builder.core.convert_pdf_to_images")
    def test_pdf_folder(self, mock_cv, mock_sz, mock_build, tmp_path):
        (tmp_path / "doc.pdf").touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        page.touch()
        mock_cv.return_value = [page]

        process_folder(tmp_path, recursive=False, dpi=150, quiet=False)

        mock_cv.assert_called_once()
        mock_build.assert_called_once()
        assert not pages_dir.exists()  # cleaned up by finally

    @pytest.mark.integration
    def test_image_folder(self, tmp_path):
        _make_image(tmp_path / "a.png", 100, 150)
        process_folder(tmp_path, recursive=False, dpi=150, quiet=False)
        assert (tmp_path / f"{tmp_path.name}.pptx").exists()

    @pytest.mark.integration
    def test_portrait_images_no_forced_landscape(self, tmp_path):
        _make_image(tmp_path / "p.png", 100, 200)
        process_folder(tmp_path, recursive=False, dpi=150, quiet=False)
        from pptx import Presentation

        prs = Presentation(str(tmp_path / f"{tmp_path.name}.pptx"))
        w = float(prs.slide_width) / 914400
        h = float(prs.slide_height) / 914400
        assert h > w  # portrait preserved

    @patch("pptx_builder.core.build_presentation")
    @patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
    @patch("pptx_builder.core.convert_pdf_to_images")
    def test_mixed_folder_prioritizes_pdfs(self, mock_cv, mock_sz, mock_build, tmp_path):
        (tmp_path / "doc.pdf").touch()
        _make_image(tmp_path / "img.png")
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        page.touch()
        mock_cv.return_value = [page]

        process_folder(tmp_path, recursive=False, dpi=150, quiet=True)
        mock_cv.assert_called_once()

    def test_empty_folder_quiet(self, tmp_path, capsys):
        process_folder(tmp_path, recursive=False, dpi=150, quiet=True)
        assert capsys.readouterr().out == ""

    def test_empty_folder_not_quiet(self, tmp_path, capsys):
        process_folder(tmp_path, recursive=False, dpi=150, quiet=False)
        assert "(empty)" in capsys.readouterr().out

    @patch("pptx_builder.core.build_presentation")
    @patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
    @patch("pptx_builder.core.convert_pdf_to_images")
    def test_recursive_visits_each_dir_once(self, mock_cv, mock_sz, mock_build, tmp_path):
        sub = tmp_path / "sub"
        sub.mkdir()
        (sub / "doc.pdf").touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)  # real image so PIL can read dimensions if needed
        mock_cv.return_value = [page]

        process_folder(tmp_path, recursive=True, dpi=150, quiet=True)
        assert mock_cv.call_count == 1  # not double-visited

    def test_image_folder_empty_after_list_images(self, tmp_path):
        """Line 406: continue when list_images() finds no files (e.g. dir with image extension)."""
        fake = tmp_path / "trick.png"
        fake.mkdir()  # directory named trick.png fools outer check but not list_images
        process_folder(tmp_path, recursive=False, dpi=150, quiet=True)
        # should complete without error


# ─── main() entrypoint ────────────────────────────────────────────────────────


class TestMainEntrypoint:
    # ── exit-code fix: --output + multiple --input ────────────────────────────
    def test_output_with_multiple_inputs_exits_1(self):
        with patch("sys.argv", ["pb", "-i", "a.pdf", "b.pdf", "-o", "out.pptx"]):
            with pytest.raises(SystemExit) as exc:
                main()
        assert exc.value.code == 1

    # ── verbose flag covers logging.DEBUG branch ───────────────────────────────
    @pytest.mark.integration
    def test_verbose_flag(self, tmp_path):
        pdf = tmp_path / "f.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pd"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "--verbose"]))
            # path, slide size, fit mode, format (JPEG), quality (accept default)
            stack.enter_context(patch("builtins.input", side_effect=[str(pdf), "1", "1", "1", ""]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()

    # ── interactive: PDF, no --output ─────────────────────────────────────────
    @pytest.mark.integration
    def test_interactive_pdf(self, tmp_path, capsys):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            # path, slide size, fit mode, format (JPEG), quality (accept default)
            stack.enter_context(patch("builtins.input", side_effect=[str(pdf), "1", "1", "1", ""]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()
        assert "saved" in capsys.readouterr().out.lower()

    # ── interactive: PDF + --output flag discards prompt ──────────────────────
    @pytest.mark.integration
    def test_interactive_pdf_with_output_flag(self, tmp_path, capsys):
        pdf = tmp_path / "t.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        out_name = str(tmp_path / "custom")
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "--output", out_name]))
            # path, slide size, fit mode, format (JPEG), quality (accept default)
            stack.enter_context(patch("builtins.input", side_effect=[str(pdf), "1", "1", "1", ""]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()
        assert "saved" in capsys.readouterr().out.lower()

    # ── interactive: folder input ─────────────────────────────────────────────
    @pytest.mark.integration
    def test_interactive_folder(self, tmp_path, capsys):
        _make_image(tmp_path / "a.png")
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            stack.enter_context(
                patch("builtins.input", side_effect=[str(tmp_path), "slides", "2", "1"])
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()
        assert "saved" in capsys.readouterr().out.lower()

    # ── interactive: unknown path → sys.exit(1) ───────────────────────────────
    def test_interactive_unknown_exits(self, tmp_path):
        txt = tmp_path / "f.txt"
        txt.touch()
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            stack.enter_context(patch("builtins.input", return_value=str(txt)))
            with pytest.raises(SystemExit) as exc:
                main()
        assert exc.value.code == 1

    # ── interactive: empty image folder → sys.exit(1) ────────────────────────
    def test_interactive_empty_image_folder_exits(self, tmp_path):
        sub = tmp_path / "sub"
        sub.mkdir()
        # Patch detect_input_type to return "folder" so we enter the folder branch,
        # then patch list_images to return [] so the "No images" exit fires.
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            stack.enter_context(patch("builtins.input", side_effect=[str(sub), "", "2", "1"]))
            stack.enter_context(patch("pptx_builder.core.detect_input_type", return_value="folder"))
            stack.enter_context(patch("pptx_builder.core.list_images", return_value=[]))
            with pytest.raises(SystemExit) as exc:
                main()
        assert exc.value.code == 1

    # ── interactive: cleanup rmtree exception is swallowed ────────────────────
    def test_interactive_cleanup_exception_is_silent(self, tmp_path):
        pdf = tmp_path / "f.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pd"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            # path, slide size, fit mode, format (JPEG), quality (accept default)
            stack.enter_context(patch("builtins.input", side_effect=[str(pdf), "1", "1", "1", ""]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            stack.enter_context(patch("shutil.rmtree", side_effect=PermissionError("locked")))
            main()  # must not raise despite rmtree failing

    # ── interactive: build exception → sys.exit(1) ────────────────────────────
    def test_interactive_build_exception_exits(self, tmp_path):
        pdf = tmp_path / "f.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pd"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb"]))
            # path, slide size, fit mode, format (JPEG), quality (accept default)
            stack.enter_context(patch("builtins.input", side_effect=[str(pdf), "1", "1", "1", ""]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.core.build_presentation", side_effect=RuntimeError("oops"))
            )
            with pytest.raises(SystemExit) as exc:
                main()
        assert exc.value.code == 1

    # ── CLI: non-existent path ────────────────────────────────────────────────
    def test_cli_nonexistent_path(self, capsys):
        with patch("sys.argv", ["pb", "-i", "/no/such/file.pdf"]):
            main()
        assert "not found" in capsys.readouterr().out.lower()

    # ── CLI: unsupported file type ────────────────────────────────────────────
    def test_cli_unsupported_type(self, tmp_path, capsys):
        txt = tmp_path / "d.txt"
        txt.touch()
        with patch("sys.argv", ["pb", "-i", str(txt)]):
            main()
        assert "Unsupported" in capsys.readouterr().out

    # ── CLI: PDF success ──────────────────────────────────────────────────────
    @pytest.mark.integration
    def test_cli_pdf_success(self, tmp_path, capsys):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(pdf)]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()
        assert "saved" in capsys.readouterr().out.lower()

    # ── CLI: PDF with --output flag ────────────────────────────────────────────
    @pytest.mark.integration
    def test_cli_pdf_with_output_flag(self, tmp_path, capsys):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(pdf), "-o", "custom"]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            main()
        assert "saved" in capsys.readouterr().out.lower()

    # ── CLI: PDF, skip existing ────────────────────────────────────────────────
    @pytest.mark.integration
    def test_cli_pdf_skip_existing(self, tmp_path, capsys):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        (tmp_path / "d.pptx").touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(pdf)]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            stack.enter_context(patch("pptx_builder.core.confirm_overwrite", return_value=False))
            main()
        assert "Skipped" in capsys.readouterr().out

    # ── CLI: PDF conversion error ─────────────────────────────────────────────
    def test_cli_pdf_conversion_error(self, tmp_path, capsys):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(pdf)]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", side_effect=RuntimeError("fail"))
            )
            main()
        assert "Failed" in capsys.readouterr().out

    # ── CLI: folder mode ──────────────────────────────────────────────────────
    def test_cli_folder_mode(self, tmp_path, capsys):
        _make_image(tmp_path / "a.png")
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(tmp_path)]))
            mock_pf = stack.enter_context(patch("pptx_builder.core.process_folder"))
            main()
        mock_pf.assert_called_once()

    # ── CLI: folder exception ─────────────────────────────────────────────────
    def test_cli_folder_exception(self, tmp_path, capsys):
        _make_image(tmp_path / "a.png")
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(tmp_path)]))
            stack.enter_context(
                patch("pptx_builder.core.process_folder", side_effect=Exception("boom"))
            )
            main()
        assert "failed" in capsys.readouterr().out.lower()

    # ── CLI: cleanup rmtree exception is swallowed ────────────────────────────
    @pytest.mark.integration
    def test_cli_cleanup_exception_is_silent(self, tmp_path, capsys):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(pdf)]))
            stack.enter_context(
                patch("pptx_builder.core.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.core.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            stack.enter_context(patch("pptx_builder.core.build_presentation"))
            stack.enter_context(patch("shutil.rmtree", side_effect=PermissionError("locked")))
            main()  # must not raise

    # ── CLI: --quiet suppresses completion message ────────────────────────────
    def test_cli_quiet_no_completion_message(self, tmp_path, capsys):
        _make_image(tmp_path / "a.png")
        with ExitStack() as stack:
            stack.enter_context(patch("sys.argv", ["pb", "-i", str(tmp_path), "--quiet"]))
            stack.enter_context(patch("pptx_builder.core.process_folder"))
            main()
        assert "complete" not in capsys.readouterr().out.lower()


# ─── web.py cleanup ───────────────────────────────────────────────────────────


class TestWebCleanup:
    def test_cleanup_temp_files_removes_dirs(self, tmp_path):
        d = Path(tempfile.mkdtemp(prefix="pptx_builder_", dir=tmp_path))
        web_module.TEMP_DIRS.append(d)
        cleanup_temp_files()
        assert not d.exists()
        assert len(web_module.TEMP_DIRS) == 0

    def test_cleanup_temp_files_tolerates_missing(self):
        web_module.TEMP_DIRS.append(Path("/tmp/nonexistent_pptx_xyz"))
        cleanup_temp_files()
        assert len(web_module.TEMP_DIRS) == 0

    def test_cleanup_old_removes_stale_builder_dirs(self, tmp_path):
        import os

        old = Path(tempfile.mkdtemp(prefix="pptx_builder_", dir=tmp_path))
        os.utime(old, (time.time() - 7200, time.time() - 7200))
        with patch("pptx_builder.web.tempfile.gettempdir", return_value=str(tmp_path)):
            cleanup_old_files()
        assert not old.exists()

    def test_cleanup_old_keeps_recent_dirs(self, tmp_path):
        recent = Path(tempfile.mkdtemp(prefix="pptx_builder_", dir=tmp_path))
        with patch("pptx_builder.web.tempfile.gettempdir", return_value=str(tmp_path)):
            cleanup_old_files()
        assert recent.exists()
        recent.rmdir()

    def test_cleanup_old_removes_stale_pdf_dirs(self, tmp_path):
        import os

        old = Path(tempfile.mkdtemp(prefix="pptx_pdf_", dir=tmp_path))
        os.utime(old, (time.time() - 7200, time.time() - 7200))
        with patch("pptx_builder.web.tempfile.gettempdir", return_value=str(tmp_path)):
            cleanup_old_files()
        assert not old.exists()

    def test_cleanup_old_prunes_dead_temp_dirs_entries(self, tmp_path):
        dead = tmp_path / "pptx_builder_dead"
        web_module.TEMP_DIRS.append(dead)
        with patch("pptx_builder.web.tempfile.gettempdir", return_value=str(tmp_path)):
            cleanup_old_files()
        assert dead not in web_module.TEMP_DIRS


# ─── web.py process_files ─────────────────────────────────────────────────────


class TestWebProcessFiles:
    def _img(self, tmp_path, name="img.png", w=100, h=100):
        return str(_make_image(tmp_path / name, w, h))

    def test_empty_files_returns_none(self):
        assert process_files([], "16:9 (Widescreen)", "Fit whole image") is None

    def test_too_many_files_raises(self, tmp_path):
        files = [self._img(tmp_path, f"i{n}.png") for n in range(MAX_FILES + 1)]
        with pytest.raises(gr.Error, match="Too many"):
            process_files(files, "16:9 (Widescreen)", "Fit whole image")

    def test_file_too_large_raises(self, tmp_path):
        img = tmp_path / "big.png"
        _make_image(img)
        # Patch st_size on the stat result
        real_stat = img.stat()
        mock_stat = MagicMock()
        mock_stat.st_size = MAX_FILE_SIZE + 1
        with patch.object(Path, "stat", return_value=mock_stat):
            with pytest.raises(gr.Error, match="too large"):
                process_files([str(img)], "16:9 (Widescreen)", "Fit whole image")

    @pytest.mark.integration
    def test_single_image_returns_pptx(self, tmp_path):
        result = process_files([self._img(tmp_path)], "16:9 (Widescreen)", "Fit whole image")
        assert result is not None and result.endswith(".pptx")
        assert Path(result).exists()

    @pytest.mark.integration
    def test_multiple_images_sorted(self, tmp_path):
        imgs = [self._img(tmp_path, f"z_{n}.png") for n in range(3)]
        result = process_files(imgs, "4:3 (Standard)", "Crop to fill")
        assert result is not None

    @pytest.mark.integration
    def test_custom_output_name_no_ext(self, tmp_path):
        result = process_files(
            [self._img(tmp_path)], "16:9 (Widescreen)", "Fit whole image", output_name="my_deck"
        )
        assert Path(result).name == "my_deck.pptx"

    @pytest.mark.integration
    def test_custom_output_name_with_ext(self, tmp_path):
        result = process_files(
            [self._img(tmp_path)], "16:9 (Widescreen)", "Fit whole image", output_name="deck.pptx"
        )
        assert Path(result).name == "deck.pptx"

    @pytest.mark.integration
    def test_multiple_files_default_name(self, tmp_path):
        imgs = [self._img(tmp_path, f"f{n}.png") for n in range(2)]
        result = process_files(imgs, "16:9 (Widescreen)", "Fit whole image")
        assert Path(result).name == "presentation.pptx"

    @pytest.mark.integration
    def test_path_traversal_stripped(self, tmp_path):
        result = process_files(
            [self._img(tmp_path)], "16:9 (Widescreen)", "Fit whole image", output_name="../escape"
        )
        assert "pptx_builder_" in Path(result).parent.name

    @pytest.mark.integration
    def test_absolute_path_in_name_stripped(self, tmp_path):
        result = process_files(
            [self._img(tmp_path)], "16:9 (Widescreen)", "Fit whole image", output_name="/etc/evil"
        )
        assert Path(result).parent.name.startswith("pptx_builder_")

    @pytest.mark.integration
    def test_fill_mode(self, tmp_path):
        result = process_files([self._img(tmp_path)], "16:9 (Widescreen)", "Crop to fill")
        assert result is not None

    def test_pdf_temp_dir_tracked_in_temp_dirs(self, tmp_path):
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "page.png"
        _make_image(page)
        pdf_path = str(tmp_path / "doc.pdf")
        Path(pdf_path).touch()

        with ExitStack() as stack:
            stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(patch("pptx_builder.web.build_presentation"))
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            process_files([pdf_path], "16:9 (Widescreen)", "Fit whole image")

        assert pages_dir in web_module.TEMP_DIRS

    def test_mixed_pdf_and_images_upload_order_preserved(self, tmp_path):
        img_a = self._img(tmp_path, "z_last.png")
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "page.png"
        _make_image(page)
        pdf_path = str(tmp_path / "doc.pdf")
        Path(pdf_path).touch()

        captured = []

        def mock_build(images, **kw):
            captured.extend(images)

        with ExitStack() as stack:
            stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.web.build_presentation", side_effect=mock_build)
            )
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            process_files([pdf_path, img_a], "16:9 (Widescreen)", "Fit whole image")

        assert captured[0] == page  # PDF page first, not z_last.png

    def test_exception_raises_gr_error(self, tmp_path):
        img = self._img(tmp_path)
        with patch("pptx_builder.web.build_presentation", side_effect=ValueError("bad")):
            with pytest.raises(gr.Error, match="bad"):
                process_files([img], "16:9 (Widescreen)", "Fit whole image")

    def test_no_images_after_pdf_returns_none(self, tmp_path):
        pdf_path = str(tmp_path / "empty.pdf")
        Path(pdf_path).touch()
        with ExitStack() as stack:
            stack.enter_context(patch("pptx_builder.web.convert_pdf_to_images", return_value=[]))
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(10.0, 7.5))
            )
            result = process_files([pdf_path], "16:9 (Widescreen)", "Fit whole image")
        assert result is None

    @pytest.mark.integration
    def test_single_pdf_auto_detects_dimensions(self, tmp_path):
        pages_dir = tmp_path / "pages"
        pages_dir.mkdir()
        page = pages_dir / "page.png"
        _make_image(page, 800, 600)
        pdf_path = str(tmp_path / "doc.pdf")
        Path(pdf_path).touch()

        with ExitStack() as stack:
            stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            mock_build = stack.enter_context(patch("pptx_builder.web.build_presentation"))
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(11.0, 8.5))
            )
            process_files([pdf_path], "16:9 (Widescreen)", "Fit whole image")

        _, kw = mock_build.call_args
        assert kw.get("slide_width_in") == 11.0 or mock_build.call_args[0][2] == 11.0


# ─── web.py launch_app ────────────────────────────────────────────────────────


class TestWebOutputFormat:
    """Web UI encoding controls reach the pipeline (ho-01 AT-02)."""

    def test_process_files_passes_format_and_quality(self, tmp_path):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        page = tmp_path / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            mock_convert = stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(11.0, 8.5))
            )
            stack.enter_context(patch("pptx_builder.web.build_presentation"))
            process_files(
                [str(pdf)],
                "16:9 (Widescreen)",
                "Fit whole image",
                dpi=200,
                out_format="png",
                quality=70,
            )
        kwargs = mock_convert.call_args.kwargs
        assert kwargs["fmt"] == "png"
        assert kwargs["quality"] == 70

    def test_process_files_clamps_quality_to_ceiling(self, tmp_path):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        page = tmp_path / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            stack.enter_context(patch("pptx_builder.web.MAX_QUALITY", 80))
            mock_convert = stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(11.0, 8.5))
            )
            stack.enter_context(patch("pptx_builder.web.build_presentation"))
            process_files(
                [str(pdf)],
                "16:9 (Widescreen)",
                "Fit whole image",
                out_format="jpeg",
                quality=100,
            )
        assert mock_convert.call_args.kwargs["quality"] == 80

    def test_process_files_defaults_match_core(self, tmp_path):
        pdf = tmp_path / "d.pdf"
        pdf.touch()
        page = tmp_path / "p.png"
        _make_image(page)
        with ExitStack() as stack:
            mock_convert = stack.enter_context(
                patch("pptx_builder.web.convert_pdf_to_images", return_value=[page])
            )
            stack.enter_context(
                patch("pptx_builder.web.pdf_first_page_size_inches", return_value=(11.0, 8.5))
            )
            stack.enter_context(patch("pptx_builder.web.build_presentation"))
            process_files([str(pdf)], "16:9 (Widescreen)", "Fit whole image")
        kwargs = mock_convert.call_args.kwargs
        assert kwargs["fmt"] == DEFAULT_FORMAT
        assert kwargs["quality"] == min(DEFAULT_QUALITY, MAX_QUALITY)

    def test_quality_visibility_follows_format(self):
        assert _quality_visibility("jpeg")["visible"] is True
        assert _quality_visibility("png")["visible"] is False

    def test_info_strip_shows_quality_cap_when_set(self):
        with patch("pptx_builder.web.MAX_QUALITY", 85):
            assert "quality" in _build_info_strip()

    def test_info_strip_omits_quality_cap_when_uncapped(self):
        with patch("pptx_builder.web.MAX_QUALITY", 100):
            assert "quality&nbsp;" not in _build_info_strip()


class TestVersion:
    def test_version_is_0_3_0(self):
        import pptx_builder

        assert pptx_builder.__version__ == "0.3.0"


class TestWebLaunchApp:
    def test_launch_calls_app_launch(self):
        from pptx_builder.web import launch_app, app

        with patch.object(app, "launch") as mock_launch:
            launch_app()
        mock_launch.assert_called_once()
        kw = mock_launch.call_args[1]
        assert kw.get("server_name") == "0.0.0.0"


# ─── cli module ───────────────────────────────────────────────────────────────


class TestCLIModule:
    def test_cli_main_is_core_main(self):
        from pptx_builder import cli
        from pptx_builder.core import main as core_main

        assert cli.main is core_main


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
